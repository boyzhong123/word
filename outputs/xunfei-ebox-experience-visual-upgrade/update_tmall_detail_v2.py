from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BASE = Path("/Users/zhong/Desktop/Cursor仓库/ChivoxMCP官网/outputs/xunfei-ebox-experience-visual-upgrade")
SRC = BASE / "xunfei-ebox-pro-experience-report-leadership-final.pptx"
OUT = BASE / "xunfei-ebox-pro-experience-report-leadership-final-v2.pptx"

IMG_DIR = Path("/Users/zhong/.codex/generated_images/019df8ae-70da-7311-a3e1-64af414e01a4")
IMG_MODULES = IMG_DIR / "ig_04e984a91b1550f90169fae63639c881909a9d6eea0fecd371.png"
IMG_JOURNEY = IMG_DIR / "ig_04e984a91b1550f90169fae7bdd3fc8190be79eb193195c62d.png"
IMG_TEST = IMG_DIR / "ig_04e984a91b1550f90169fae80b208c8190afe29063967230cf.png"
IMG_IMPACT = IMG_DIR / "ig_04e984a91b1550f90169fae860d6dc81909a931810f0e91d25.png"
IMG_REAL_HERO = BASE.parent / "tmall-scan/images2/O1CN01A9yY3R2EOglYYD5iO_!!2813588735.png"

FONT = "Microsoft YaHei"
SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)


def c(hex_color: str) -> RGBColor:
    hex_color = hex_color.strip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


COLORS = {
    "ink": c("0B172A"),
    "muted": c("566070"),
    "soft": c("F7F5F1"),
    "panel": c("FFFFFF"),
    "orange": c("FF5A1F"),
    "orange2": c("FF8A2A"),
    "blue": c("256DFF"),
    "teal": c("0EB8A8"),
    "green": c("05B26B"),
    "red": c("F04444"),
    "yellow": c("F2A100"),
    "line": c("E4E7E2"),
    "lav": c("F0F3FF"),
    "peach": c("FFF1E8"),
    "mint": c("EAF8F1"),
    "sky": c("EDF5FF"),
    "rose": c("FFF0F0"),
}


def clear(slide):
    for shape in list(slide.shapes):
        slide.shapes._spTree.remove(shape._element)


def set_fill(shape, fill_color, transparency=0):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_color
    if transparency:
        fill.transparency = transparency


def set_line(shape, line_color=None, width=1, transparency=0):
    line = shape.line
    if line_color is None:
        line.fill.background()
        return
    line.color.rgb = line_color
    line.width = Pt(width)
    if transparency:
        line.transparency = transparency


def add_text(slide, x, y, w, h, text, size=20, color=None, bold=False, align="left", weight=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    lines = str(text).split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color or COLORS["ink"]
    return box


def add_rich_line(slide, x, y, w, h, pieces, size=20, align="left"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    for text, color, bold in pieces:
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_bg(slide, tint="soft"):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    set_fill(shape, COLORS[tint])
    set_line(shape, None)
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)


def add_title(slide, idx, title, subtitle=None, eyebrow="EBOX PRO EXPERIENCE"):
    add_text(slide, 0.72, 0.46, 2.2, 0.25, eyebrow, 9.5, COLORS["orange"], True)
    add_text(slide, 11.82, 0.43, 0.7, 0.25, f"{idx:02d}", 10, COLORS["muted"], False, "right")
    add_text(slide, 0.72, 0.86, 8.9, 0.58, title, 25, COLORS["ink"], True)
    if subtitle:
        add_text(slide, 0.72, 1.48, 10.8, 0.38, subtitle, 12.5, COLORS["muted"])
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(1.92), Inches(11.9), Inches(0.01))
    set_fill(line, COLORS["line"])
    set_line(line, None)


def add_footer(slide, text):
    add_text(slide, 0.72, 7.12, 11.7, 0.22, text, 7.2, c("8A8F99"))


def add_card(slide, x, y, w, h, title, body, accent="orange", fill="panel", title_size=15, body_size=10.5):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(box, COLORS[fill])
    set_line(box, COLORS["line"], 0.7)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.055), Inches(h))
    set_fill(strip, COLORS[accent])
    set_line(strip, None)
    add_text(slide, x + 0.18, y + 0.18, w - 0.34, 0.28, title, title_size, COLORS[accent], True)
    add_text(slide, x + 0.18, y + 0.58, w - 0.34, h - 0.72, body, body_size, COLORS["muted"])
    return box


def add_metric(slide, x, y, w, h, num, label, accent="orange", sub=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, COLORS["panel"])
    set_line(shape, COLORS["line"], 0.8)
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.08))
    set_fill(top, COLORS[accent])
    set_line(top, None)
    add_text(slide, x + 0.18, y + 0.22, w - 0.36, 0.34, num, 24, COLORS[accent], True)
    add_text(slide, x + 0.18, y + 0.67, w - 0.36, 0.24, label, 10.5, COLORS["ink"], True)
    if sub:
        add_text(slide, x + 0.18, y + 0.97, w - 0.36, 0.22, sub, 8.2, COLORS["muted"])
    return shape


def add_pill(slide, x, y, w, h, text, accent="orange", fill=None, size=10.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shp, COLORS[fill] if fill else COLORS[accent])
    set_line(shp, None)
    add_text(slide, x, y + 0.06, w, h - 0.07, text, size, c("FFFFFF") if not fill else COLORS[accent], True, "center")
    return shp


def add_icon(slide, x, y, text, accent="orange", size=17, fill=None):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.52), Inches(0.52))
    set_fill(circle, COLORS[fill] if fill else COLORS["panel"])
    set_line(circle, COLORS[accent], 1.4)
    add_text(slide, x, y + 0.09, 0.52, 0.3, text, size, COLORS[accent], True, "center")
    return circle


def add_picture(slide, path, x, y, w, h=None, transparency_panel=False):
    if not path.exists():
        return None
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    if h:
        pic.height = Inches(h)
    if transparency_panel:
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x - 0.03), Inches(y - 0.03), Inches(w + 0.06), Inches((h or pic.height / 914400) + 0.06))
        set_fill(panel, COLORS["panel"], transparency=8)
        set_line(panel, COLORS["line"], 0.5)
        slide.shapes._spTree.remove(panel._element)
        slide.shapes._spTree.insert(slide.shapes._spTree.index(pic._element), panel._element)
    return pic


def rebuild_cover(slide):
    clear(slide)
    add_bg(slide, "soft")
    add_text(slide, 0.78, 0.52, 2.8, 0.32, "COMPETITOR EXPERIENCE REPORT", 9.5, COLORS["orange"], True)
    add_rich_line(
        slide, 0.78, 1.15, 8.9, 0.72,
        [("科大讯飞 AI英语宝 ", COLORS["ink"], True), ("EBOX Pro", COLORS["orange"], True)],
        31,
    )
    add_text(slide, 0.78, 1.96, 8.0, 0.46, "硬件产品体验报告｜基于真机验证方向与天猫详情页信息补充", 15, COLORS["muted"])
    add_picture(slide, IMG_REAL_HERO, 9.18, 0.48, 2.72, 5.12, True)
    add_metric(slide, 0.78, 3.25, 1.78, 1.18, "6合1", "家庭端定位", "orange", "六类能力合并")
    add_metric(slide, 2.78, 3.25, 1.78, 1.18, "近400", "口语话题", "teal", "覆盖新课标场景")
    add_metric(slide, 4.78, 3.25, 1.78, 1.18, "10000+", "中高考真题", "blue", "详情页口径")
    add_metric(slide, 6.78, 3.25, 1.78, 1.18, "95%+", "教材覆盖率", "green", "详情页口径")
    add_metric(slide, 8.78, 3.25, 1.78, 1.18, "105g", "便携硬件", "yellow", "3英寸屏/1800mAh")
    add_card(slide, 0.78, 5.05, 5.48, 1.05, "一句话判断", "EBOX Pro 的核心不是单点硬件，而是把英语练习、AI反馈、内容资源与家长安心包装成低价家庭端入口。", "orange", "panel", 15, 12)
    add_card(slide, 6.55, 5.05, 5.48, 1.05, "对驰声的关键意义", "短期更偏 2C 家庭自练；若讯飞把教师端、成套题和学校渠道补齐，将挤压课堂听说练测硬件入口预期。", "red", "panel", 15, 12)
    add_footer(slide, "信息来源：讯飞官网/公开报道/天猫详情页截图/驰声公开资料与现有产品材料；需以已购真机复核为准。")


def rebuild_exec(slide):
    clear(slide)
    add_bg(slide, "soft")
    add_title(slide, 2, "管理层速读：天猫详情页把它包装成“家庭英语闭环入口”",
              "这版补充重点：不是再介绍一个复读机，而是看清它正在争夺哪些家庭学习场景。")
    add_card(slide, 0.82, 2.25, 2.75, 2.05, "1. 入口定位", "从“听力宝/单词机”升级为背词、口语、听力、写作、真题与家长管控的组合型终端。", "orange", "panel", 14, 11.2)
    add_card(slide, 3.84, 2.25, 2.75, 2.05, "2. 内容资产", "强调同步教材、时文、名著、真题、近400口语话题、18个单词训练营，核心是持续内容更新。", "blue", "panel", 14, 11.2)
    add_card(slide, 6.86, 2.25, 2.75, 2.05, "3. 家长理由", "0游戏、可通话定位支付、远程管控、学情可视，让家长把它视为学习机+手表替代。", "green", "panel", 14, 11.2)
    add_card(slide, 9.88, 2.25, 2.75, 2.05, "4. 竞品边界", "目前仍是家庭个人练习，不是课堂统一听说练测系统；但升级到学校端会带来硬件入口冲击。", "red", "panel", 14, 11.2)
    add_pill(slide, 0.88, 4.84, 1.18, 0.34, "结论", "orange", size=10)
    add_text(slide, 2.25, 4.73, 9.7, 0.68, "它值得持续跟踪，因为讯飞正在把“语音评测能力”放进家庭高频练习场景，再用硬件低价和家长安心降低购买门槛。", 17, COLORS["ink"], True)
    add_footer(slide, "补充依据：用户提供的天猫详情页截图，包含功能、内容、家长管控、参数与竞品对比模块。")


def rebuild_definition(slide):
    clear(slide)
    add_bg(slide, "soft")
    add_title(slide, 3, "产品定义：6合1不是堆功能，而是合并家长会付费的场景",
              "详情页核心表达是“省钱更省心”：学习、陪练、通讯、安全、内容订阅被打包进一个小屏硬件。")
    add_picture(slide, IMG_MODULES, 7.85, 1.38, 4.28, 3.05, True)
    modules = [
        ("听力机", "熏听、复读、跟读评分\n校内教材与时文精讲", "teal"),
        ("口语私教", "中英混合输入\n24小时主动引导", "blue"),
        ("电话手表", "通话、定位、支付\n上课模式与微聊管控", "green"),
        ("单词机", "诊学测、训练营\n故事/PK/真题检测", "orange"),
        ("电子词典", "查词、翻译、自建词书\n拍照一键导入", "yellow"),
        ("写作顾问", "作文批改、润色建议\n结构与整体点评", "red"),
    ]
    for i, (t, b, accent) in enumerate(modules):
        x = 0.82 + (i % 2) * 3.28
        y = 2.12 + (i // 2) * 1.38
        add_card(slide, x, y, 2.9, 1.06, t, b, accent, "panel", 13.5, 9.3)
    add_text(slide, 0.82, 6.45, 11.4, 0.34, "判断：详情页在刻意摆脱“单词机/复读机”的窄认知，把它塑造成家庭英语学习、通信与管控的低价硬件入口。", 12.5, COLORS["ink"], True)
    add_footer(slide, "来自天猫截图：升级6合1、电话手表平替、0游戏、没网也能用、比橙子轻等卖点。")


def rebuild_function_map(slide):
    clear(slide)
    add_bg(slide, "soft")
    add_title(slide, 7, "功能地图：详情页呈现的是“学-练-测-管”的闭环",
              "把每个功能放回体验链路看，核心不是功能数量，而是一次练习后能否继续反馈和复练。")
    cards = [
        ("Aa", "背单词", "AI诊学测、训练营、真人原音、故事、PK", "orange", "peach"),
        ("听", "听力", "同步教材、时文名著真题、老师精讲、高频词", "teal", "mint"),
        ("说", "口语", "中英混合输入、近400话题、即时引导", "blue", "sky"),
        ("测", "真题", "5000+历年真题、考点解析、个性化推送", "green", "mint"),
        ("写", "写作", "作文批改、AI润色、查词翻译、听写", "red", "rose"),
        ("管", "家长", "通话定位支付、学情可视、时长与权限管控", "yellow", "peach"),
    ]
    for i, (icon, title, body, accent, fill) in enumerate(cards):
        x = 0.86 + (i % 3) * 4.08
        y = 2.26 + (i // 3) * 1.9
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.48), Inches(1.38))
        set_fill(shp, COLORS[fill])
        set_line(shp, COLORS["line"], 0.7)
        add_icon(slide, x + 0.28, y + 0.32, icon, accent, 15, "panel")
        add_text(slide, x + 0.98, y + 0.28, 2.08, 0.3, title, 15.5, COLORS[accent], True)
        add_text(slide, x + 0.98, y + 0.72, 2.22, 0.4, body, 9.2, COLORS["muted"])
    add_pill(slide, 0.86, 6.25, 1.28, 0.32, "体验抓手", "orange", size=9.5)
    add_text(slide, 2.32, 6.19, 9.8, 0.35, "后续真机验证要看：孩子是否愿意完成第二轮复练，家长是否持续查看学情并设置管控。", 12.5, COLORS["ink"], True)
    add_footer(slide, "补充自天猫详情页：单词、听力、口语、真题、作文、家长管控等功能模块。")


def rebuild_listening(slide):
    clear(slide)
    add_bg(slide, "soft")
    add_title(slide, 9, "听力与内容：从“能听”转向“听得懂、有重点、能复练”",
              "详情页强调教研精讲、课本同步、时文名著真题与高频词提取，目的在于降低理解门槛。")
    items = [
        ("同步教材", "覆盖全国主流教材版本，新课标同步更新，服务校内预复习。", "orange"),
        ("时文/名著/真题", "News in Levels、Kids News、名著与真题精讲，每日更新新鲜实事。", "blue"),
        ("老师精讲", "每篇听力配教研老师精讲，搭建理解脚手架。", "teal"),
        ("高频词提取", "从考试文章中提取小学/KET/PET/中高考等高频词。", "green"),
    ]
    for i, (t, b, accent) in enumerate(items):
        add_card(slide, 0.82, 2.16 + i * 1.02, 4.55, 0.78, t, b, accent, "panel", 12.5, 9.2)
    # Right-side workflow
    steps = [("内容输入", "教材/时文/名著/真题", "orange"), ("结构精讲", "词汇/语法/句型", "blue"), ("重点提取", "高频词/考点", "teal"), ("复练反馈", "跟读/测验/报告", "green")]
    for i, (t, b, accent) in enumerate(steps):
        x = 6.25 + (i % 2) * 2.85
        y = 2.28 + (i // 2) * 1.78
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.42), Inches(1.16))
        set_fill(shp, COLORS["panel"])
        set_line(shp, COLORS[accent], 1.1)
        add_text(slide, x + 0.22, y + 0.22, 1.95, 0.25, t, 14, COLORS[accent], True)
        add_text(slide, x + 0.22, y + 0.58, 2.0, 0.28, b, 9.6, COLORS["muted"])
    add_rich_line(slide, 6.25, 5.95, 5.35, 0.45, [("体验假设：", COLORS["orange"], True), ("如果精讲质量足够稳定，它会把“听力机”拉升为轻量教辅。", COLORS["ink"], True)], 12.5)
    add_footer(slide, "截图依据：同步教材、时文+名著+真题精讲、资深教研团队细致精讲、考试高频词提取。")


def rebuild_words(slide):
    clear(slide)
    add_bg(slide, "soft")
    add_title(slide, 10, "单词与写作：用“诊-学-练-测-用”弱化孤立背词",
              "它把单词学习做成一条可复现的路径：先摸底，再按内容和考试资源复练。")
    steps = [("诊", "AI诊断摸底"), ("学", "真人原音/考点"), ("练", "训练营/速听"), ("测", "真题/听写"), ("用", "故事/作文/口语")]
    for i, (t, b) in enumerate(steps):
        x = 0.9 + i * 2.36
        add_icon(slide, x, 2.38, t, ["orange", "blue", "teal", "green", "red"][i], 16, "panel")
        add_text(slide, x - 0.24, 3.02, 1.0, 0.24, b, 9.5, COLORS["ink"], True, "center")
        if i < 4:
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.62), Inches(2.64), Inches(1.52), Inches(0.04))
            set_fill(bar, COLORS["line"])
            set_line(bar, None)
    cards = [
        ("18个免费单词训练营", "小学/初中/高中分层攻克核心词汇", "orange"),
        ("单词真人原音讲解", "常见语块、同义辨析、高频考点", "blue"),
        ("单词秒变故事书", "星火×DeepSeek，把生词放入语境", "teal"),
        ("PK竞技模式", "把手机式即时反馈转化为学习动机", "green"),
        ("作文批改与润色", "结构、精细批改、整体点评", "red"),
        ("查词翻译/听写", "写作业时顺手练发音与听写", "yellow"),
    ]
    for i, (t, b, accent) in enumerate(cards):
        x = 0.92 + (i % 3) * 4.02
        y = 4.05 + (i // 3) * 0.95
        add_card(slide, x, y, 3.42, 0.74, t, b, accent, "panel", 11.2, 8.3)
    add_footer(slide, "截图依据：AI诊学测、五步闭环、18个单词训练营、真人原音讲解、单词故事、PK、作文批改与听写。")


def rebuild_content(slide):
    clear(slide)
    add_bg(slide, "soft")
    add_title(slide, 11, "内容资产：详情页把“资源广度”作为核心护城河表达",
              "这些数字是营销口径，需要真机抽样验证可用性、更新频率和题目颗粒度。")
    metrics = [
        ("21年", "教育资源沉淀", "orange", "启蒙到大学全覆盖"),
        ("10000+", "中高考真题", "red", "详情页口径"),
        ("95%+", "教材版本覆盖率", "green", "详情页口径"),
        ("近400", "口语话题", "blue", "覆盖新课标"),
        ("5000+", "历年真题", "teal", "结合学情推送"),
        ("18个", "免费单词训练营", "yellow", "分学段词汇"),
    ]
    for i, (num, label, accent, sub) in enumerate(metrics):
        add_metric(slide, 0.82 + (i % 3) * 4.05, 2.2 + (i // 3) * 1.45, 3.35, 1.08, num, label, accent, sub)
    add_card(slide, 0.82, 5.48, 5.65, 0.82, "对2C体验的价值", "内容广度越强，越容易让家长相信“长期可用”，并降低硬件一次性购买犹豫。", "orange", "panel", 12.5, 9.5)
    add_card(slide, 6.75, 5.48, 5.65, 0.82, "对课堂产品的边界", "资源多不等于课堂闭环。学校场景还需要统一任务、教师端管理、班级数据和成套听说题。", "blue", "panel", 12.5, 9.5)
    add_footer(slide, "截图依据：21年教育资源沉淀、10000+真题、95%+教材覆盖、近400口语话题、5000+历年真题、18个训练营。")


def rebuild_parent(slide):
    clear(slide)
    add_bg(slide, "soft")
    add_title(slide, 12, "家长与安全：它卖的不只是学习，也卖“可控的独立设备”",
              "电话手表平替、家长管控、0游戏、离线/SIM，是降低家长购买阻力的关键包装。")
    groups = [
        ("电话手表平替", [("通话", "随时联系"), ("定位", "安全感"), ("支付", "轻量外出")], "orange"),
        ("家长远程管控", [("学情可视", "随时查看"), ("上课模式", "时间权限"), ("微聊/通话", "关系管控")], "blue"),
        ("学习环境", [("不能下载游戏", "降低分心"), ("无娱乐干扰", "沉浸学习"), ("离线/SIM", "在校/在途/在外")], "green"),
    ]
    for gi, (title, rows, accent) in enumerate(groups):
        x = 0.82 + gi * 4.03
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.18), Inches(3.48), Inches(3.35))
        set_fill(panel, COLORS["panel"])
        set_line(panel, COLORS["line"], 0.7)
        add_text(slide, x + 0.25, 2.45, 2.6, 0.28, title, 15, COLORS[accent], True)
        for i, (a, b) in enumerate(rows):
            add_icon(slide, x + 0.28, 3.02 + i * 0.72, str(i + 1), accent, 12, "panel")
            add_text(slide, x + 0.92, 3.02 + i * 0.72, 1.5, 0.22, a, 11.2, COLORS["ink"], True)
            add_text(slide, x + 0.92, 3.28 + i * 0.72, 1.95, 0.18, b, 8.6, COLORS["muted"])
    add_pill(slide, 0.82, 6.18, 1.15, 0.32, "验证点", "red", size=9.5)
    add_text(slide, 2.15, 6.1, 9.8, 0.36, "这些能力若体验顺滑，会让家长把它当成“手机替代品”；若管控或定位体验粗糙，卖点会快速掉分。", 12.2, COLORS["ink"], True)
    add_footer(slide, "截图依据：电话手表平替、通话/定位/支付、时长/内容权限、家长远程管控、不能下载游戏、离线/SIM。")


def rebuild_channel(slide):
    clear(slide)
    add_bg(slide, "soft")
    add_title(slide, 13, "同类竞品对比：详情页在主动抢“家庭学习硬件”的心智",
              "它把自己放在听力宝、单词机、APP的对面，强调全链路、全年龄段、专注和便携。")
    col_x = [0.82, 3.4, 5.95, 8.5, 11.05]
    col_w = [2.22, 2.15, 2.15, 2.15, 1.55]
    headers = ["维度", "EBOX Pro", "听力宝类", "单词通类", "APP类"]
    for i, h in enumerate(headers):
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(col_x[i]), Inches(2.08), Inches(col_w[i]), Inches(0.48))
        set_fill(shp, COLORS["orange"] if i == 1 else COLORS["panel"])
        set_line(shp, COLORS["line"], 0.6)
        add_text(slide, col_x[i], 2.2, col_w[i], 0.22, h, 10.2, c("FFFFFF") if i == 1 else COLORS["ink"], True, "center")
    rows = [
        ("核心功能", "背单词+口语+真题\n全链路闭环", "听力强\n背词弱", "重拼写\n真题少", "手机分心"),
        ("学习阶段", "启蒙到应试\n全年龄段", "更偏大孩子", "偏娱乐化", "单一阶段"),
        ("专注度", "0游戏\n纯学习", "功能多易分散", "图新鲜易分心", "分心严重"),
        ("便携性", "105g\n可替代手机", "大且重", "黑白屏弱", "依赖手机"),
    ]
    for r, row in enumerate(rows):
        y = 2.72 + r * 0.76
        for i, txt in enumerate(row):
            fill = "peach" if i == 1 else "panel"
            shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(col_x[i]), Inches(y), Inches(col_w[i]), Inches(0.62))
            set_fill(shp, COLORS[fill])
            set_line(shp, COLORS["line"], 0.5)
            color = COLORS["orange"] if i == 1 else COLORS["muted"]
            add_text(slide, col_x[i] + 0.08, y + 0.1, col_w[i] - 0.16, 0.36, txt, 8.6 if i else 9.2, color if i == 1 else COLORS["ink"] if i == 0 else COLORS["muted"], i <= 1, "center")
    add_card(slide, 0.82, 6.05, 11.72, 0.62, "读法", "这张对比不是客观评测，而是销售心智设计：把家庭端硬件需求集中到“专注、闭环、低价、可控”四个购买理由上。", "blue", "panel", 11.5, 9.2)
    add_footer(slide, "截图依据：英语宝推荐首选与同类竞品对比表。")


def rebuild_score(slide):
    clear(slide)
    add_bg(slide, "soft")
    add_title(slide, 14, "体验评分框架：把详情页卖点转成真机验证问题",
              "领导汇报里建议少做结论式夸赞，多呈现“已看到的强项”和“必须验证的风险”。")
    left = [
        ("家庭端闭环强", "单词、听力、口语、作文、真题、家长管控组合完整。"),
        ("内容口径丰富", "同步教材、真题、话题、训练营、时文名著形成长期使用理由。"),
        ("硬件门槛低", "3英寸、105g、离线/SIM、0游戏，适合家庭和通勤场景。"),
    ]
    right = [
        ("评分稳定性", "同一学生、同一题、多次录音在不同噪声下分数是否稳定。"),
        ("复练完成率", "孩子拿到纠错后是否愿意再练一次，而不是看完就退出。"),
        ("内容可用性", "教材、真题、话题是否真能覆盖目标年级与本地考试诉求。"),
        ("家长端粘性", "家长是否持续查看学情并使用上课/微聊/通话管控。"),
    ]
    add_text(slide, 0.9, 2.18, 4.8, 0.3, "已看到的强项", 16, COLORS["green"], True)
    for i, (t, b) in enumerate(left):
        add_card(slide, 0.9, 2.68 + i * 1.0, 4.75, 0.76, t, b, "green", "panel", 11.5, 8.8)
    add_text(slide, 6.22, 2.18, 4.8, 0.3, "必须用真机验证", 16, COLORS["red"], True)
    for i, (t, b) in enumerate(right):
        add_card(slide, 6.22, 2.58 + i * 0.82, 5.78, 0.62, t, b, "red", "panel", 10.8, 8.3)
    add_footer(slide, "注：详情页卖点不能等同于真实体验，评分报告与家长端需要用统一题目、统一学生样本与统一噪声条件复测。")


def rebuild_test_plan(slide):
    clear(slide)
    add_bg(slide, "soft")
    add_title(slide, 15, "真机验证路线：围绕“孩子会不会复练、家长会不会持续看”",
              "已购真机后，建议用统一题目、统一学生样本、统一环境把营销卖点转成可判断证据。")
    add_picture(slide, IMG_TEST, 8.32, 1.62, 3.75, 2.7, True)
    tracks = [
        ("学习闭环", ["单词诊学测是否能推荐学习方案", "听力精讲后是否引导复读/跟读", "真题解析是否能回流到错题复练"], "orange"),
        ("语音评测", ["发音/语调/流畅度/语法维度是否透明", "同题多次录音分数波动", "弱项反馈是否可执行"], "blue"),
        ("家长与硬件", ["离线与SIM场景是否稳定", "上课/微聊/通话管控是否顺手", "学情报告是否让家长愿意持续查看"], "green"),
    ]
    for i, (title, bullets, accent) in enumerate(tracks):
        x = 0.82 + i * 2.55
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.24), Inches(2.22), Inches(3.45))
        set_fill(panel, COLORS["panel"])
        set_line(panel, COLORS["line"], 0.7)
        add_text(slide, x + 0.2, 2.5, 1.75, 0.25, title, 14, COLORS[accent], True)
        for j, b in enumerate(bullets):
            add_icon(slide, x + 0.2, 3.02 + j * 0.72, "✓", accent, 12, "panel")
            add_text(slide, x + 0.74, 3.0 + j * 0.72, 1.25, 0.42, b, 7.8, COLORS["muted"])
    add_card(slide, 8.34, 4.72, 3.72, 0.95, "建议输出物", "一页评分稳定性表 + 一页复练完成率漏斗 + 一页家长端使用截图证据。", "red", "panel", 12, 9.2)
    add_footer(slide, "验证对象：已购 EBOX Pro 真机；对照组可选驰声课堂硬件/手机APP/同类听力设备。")


def rebuild_impact(slide):
    clear(slide)
    add_bg(slide, "soft")
    add_title(slide, 16, "对驰声的冲击：当前是家庭端，真正风险在“学校端升级路径”",
              "这部分只是体验报告中的竞争判断，不应取代前面的产品理解。")
    add_picture(slide, IMG_IMPACT, 8.25, 1.62, 3.85, 2.7, True)
    stages = [
        ("当前形态", "2C家庭自练\n家长购买与管控", "冲击有限", "green"),
        ("可升级方向", "教师端、成套题\n硬件套装、校内试点", "需要监控", "orange"),
        ("形成威胁", "课堂统一练测评\n班级数据与学校渠道", "直接竞争", "red"),
    ]
    for i, (t, b, tag, accent) in enumerate(stages):
        x = 0.92 + i * 2.43
        add_card(slide, x, 2.35, 2.08, 2.02, t, b, accent, "panel", 13, 10.5)
        add_pill(slide, x + 0.28, 3.92, 1.2, 0.3, tag, accent, size=9)
        if i < 2:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.12), Inches(3.05), Inches(0.62), Inches(0.28))
            set_fill(arrow, COLORS["line"])
            set_line(arrow, None)
    add_card(slide, 0.92, 5.16, 6.95, 0.86, "判断", "讯飞已经具备语音评测、内容资源、低价硬件和家长信任。若进入学校课堂，它缺的不是单项能力，而是课堂组织与交付体系。", "red", "panel", 12.5, 9.2)
    add_card(slide, 8.28, 5.16, 3.8, 0.86, "我们的观察口径", "看是否出现：教师端、班级报告、学校套装、区域试点、成套听说题。", "blue", "panel", 12.5, 9.2)
    add_footer(slide, "竞争判断基于公开详情页与驰声现有课堂产品定位，不代表对方实际学校渠道计划。")


def rebuild_difference(slide):
    clear(slide)
    add_bg(slide, "soft")
    add_title(slide, 17, "与驰声产品的差异：一个是家庭闭环，一个是课堂系统",
              "领导需要看到边界：讯飞强在家庭端入口，驰声要守住学校端“组织教学与规模化评测”的价值。")
    headers = ["维度", "EBOX Pro", "驰声 AI智慧课堂/智课宝"]
    xs = [0.86, 3.05, 7.02]
    ws = [1.72, 3.48, 4.98]
    for i, h in enumerate(headers):
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(xs[i]), Inches(2.1), Inches(ws[i]), Inches(0.5))
        set_fill(shp, COLORS["orange"] if i == 1 else COLORS["blue"] if i == 2 else COLORS["panel"])
        set_line(shp, COLORS["line"], 0.5)
        add_text(slide, xs[i], 2.22, ws[i], 0.22, h, 10.5, c("FFFFFF") if i else COLORS["ink"], True, "center")
    rows = [
        ("目标用户", "家长购买，孩子个人练习", "学校/教师组织全班听说练测"),
        ("核心闭环", "个人练习-评分-复练-家长查看", "课堂任务-全员作答-AI批阅-班级讲评"),
        ("内容形态", "教材/真题/话题/训练营广覆盖", "课堂听说读写题型与教学资源沉淀"),
        ("硬件组合", "单台小屏设备，强调便携和安全", "智课宝+充电柜+路由器，面向班级部署"),
        ("数据价值", "个人学情与家长端查看", "班级数据、教师报告、课堂诊断与管理"),
    ]
    for r, row in enumerate(rows):
        y = 2.78 + r * 0.62
        for i, txt in enumerate(row):
            shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(xs[i]), Inches(y), Inches(ws[i]), Inches(0.48))
            set_fill(shp, COLORS["panel"])
            set_line(shp, COLORS["line"], 0.45)
            add_text(slide, xs[i] + 0.1, y + 0.1, ws[i] - 0.2, 0.25, txt, 8.3 if i else 8.8, COLORS["ink"] if i == 0 else COLORS["muted"], i == 0, "center")
    add_card(slide, 0.86, 6.15, 11.14, 0.58, "建议对外口径", "不要把 EBOX Pro 只说成“个人学习机”。更准确的说法是：它正在验证家庭端英语练习入口，但还不是学校课堂听说练测系统。", "orange", "panel", 10.8, 8.8)
    add_footer(slide, "驰声资料依据：用户提供的 AI智慧课堂产品 PPT 与驰声公开产品信息。")


def rebuild_final(slide):
    clear(slide)
    add_bg(slide, "soft")
    add_title(slide, 18, "建议动作：用真机证据守住课堂边界，同时监控讯飞升级信号",
              "这份报告建议作为内部跟踪基线，后续用真实体验数据更新。")
    actions = [
        ("1", "建立真机证据库", "统一题目、统一学生样本、统一噪声环境，记录评分日志、复练链路和家长端截图。", "orange"),
        ("2", "固化差异化话术", "驰声强调课堂听说练测、教师端组织、班级报告与成套题，不陷入单机参数对比。", "blue"),
        ("3", "监控学校端信号", "持续关注教师端、学校套装、成套听说题、区域试点、课堂报告等能力是否出现。", "red"),
    ]
    for i, (num, title, body, accent) in enumerate(actions):
        x = 0.9 + i * 4.03
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.33), Inches(3.48), Inches(2.45))
        set_fill(shp, COLORS["panel"])
        set_line(shp, COLORS["line"], 0.8)
        add_icon(slide, x + 0.28, 2.68, num, accent, 16, "panel")
        add_text(slide, x + 0.95, 2.68, 2.05, 0.28, title, 14.5, COLORS[accent], True)
        add_text(slide, x + 0.32, 3.28, 2.82, 0.86, body, 10, COLORS["muted"])
    add_rich_line(slide, 1.12, 5.62, 10.8, 0.44, [("最终判断：", COLORS["orange"], True), ("EBOX Pro 当前不是课堂产品替代品，但它在家庭端做到了“低价硬件 + AI反馈 + 内容资源 + 家长安心”的强包装，值得持续重点跟踪。", COLORS["ink"], True)], 14)
    add_footer(slide, "下一版建议补充：真机录音样本、评分波动表、家长App路径截图、与驰声课堂任务链路的并排对比。")


def main():
    prs = Presentation(str(SRC))
    rebuilders = {
        0: rebuild_cover,
        1: rebuild_exec,
        2: rebuild_definition,
        6: rebuild_function_map,
        8: rebuild_listening,
        9: rebuild_words,
        10: rebuild_content,
        11: rebuild_parent,
        12: rebuild_channel,
        13: rebuild_score,
        14: rebuild_test_plan,
        15: rebuild_impact,
        16: rebuild_difference,
        17: rebuild_final,
    }
    for idx, fn in rebuilders.items():
        fn(prs.slides[idx])
    prs.save(str(OUT))
    print(OUT)


if __name__ == "__main__":
    main()